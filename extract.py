from spacy.lang.en.stop_words import STOP_WORDS
from string import punctuation
from collections import Counter
from heapq import nlargest
import json
import spacy
import pandas as pd
from pypdf import PdfReader
import textract
import docx2txt
import os

from pptx import Presentation

#Preprocessing
def ReadFiles(file):

    f = file

    filename = f.split("/")[-1]

    df_corpus = ""
    fn = ""
        
    # checking if it is a PDF file
    if os.path.isfile(f) and str(filename).endswith('.pdf'):

        reader = PdfReader(f)
        doc_pages_txt = ''
    
        for page in reader.pages:
            # extract text
            doc_pages_txt = doc_pages_txt + page.extract_text()

        if doc_pages_txt != '':
            df_corpus = doc_pages_txt

        #If the above returns as False, we run the OCR library textract to #convert scanned/image based PDF files into text.
        else:
            ocr_text = textract.process(f, method='tesseract', lang='eng')
            df_corpus = ocr_text

        fn = filename.replace(".pdf", "")
            
    # checking if it is a docx file
    elif os.path.isfile(f) and str(filename).endswith('.docx'):
        # extract text
        docx_text = docx2txt.process(f)
    
        df_corpus = docx_text

        fn = filename.replace(".docx", "")
    
    # checking if it is a pptx file
    elif os.path.isfile(f) and str(filename).endswith('.pptx'):

        #extract text
        prs = Presentation(f)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)

        #Remove blank values and join list of strings into 1 full string    
        text_runs.remove(" ")
        pptx_text = ". ".join(text_runs) + "."
        pptx_text = " ".join(pptx_text.split())

        df_corpus = pptx_text

        fn = filename.replace(".pptx", "")

    return df_corpus, fn

def ReadFiles2(file):

    f = file

    filename = f.split("/")[-1]

    keywords = []
    summary = ""
    fn = ""

    # checking if it is a csv or xlsx file
    if os.path.isfile(f) and str(filename).endswith('.csv'):
        df = pd.read_csv(f)
        fn = filename.replace(".csv", "")
    else:
        df = pd.read_excel(f)
        fn = filename.replace(".xlsx", "")

    #Selects final 5 column headers as keywords or all column headers if less than 5 are present
    if len(df.columns) <= 5:
        keywords = df.columns
        temp_kw = ", ".join(df.columns)
        summary = "This document contains tabular data with column headers including: " + temp_kw + "."
    else:
        keywords = df.columns[-5:]
        temp_kw = ", ".join(df.columns[-5:])
        summary = "This document contains tabular data with column headers including: " + temp_kw + "."

    return fn, summary, list(keywords)

#Summarization and Keyword Extraction

def keyword_extract(model, text):

    nlp = model
    doc = nlp(text)

    #Gets words with specific parts of speech and finds frequency of each word
    keyword = []
    stopwords = list(STOP_WORDS)
    pos_tag = ['PROPN', 'ADJ', 'NOUN', 'VERB']
    for token in doc:
        if token.text in stopwords or token.text in punctuation:
            continue
        if token.pos_ in pos_tag:
            keyword.append(token.text)
    
    freq_word = Counter(keyword)
    freq_word.most_common(5)

    max_freq = Counter(keyword).most_common(1)[0][1]
    for word in freq_word.keys():
        freq_word[word] = (freq_word[word]/max_freq)
    freq_word.most_common(5)

    #Select top 5 most frequently occuring words
    kw_list = []
    for i in freq_word.most_common(5):
        kw_list.append(i[0])

    return freq_word, doc, kw_list

def sentence_weighting(text, keyword):
    doc = text
    freq_word = keyword
    sent_strength={}

    #Weight sentences based on how frequent previous keywords that were selected appear
    for sent in doc.sents:
        for word in sent:
            if word.text in freq_word.keys():
                if sent in sent_strength.keys():
                    sent_strength[sent]+=freq_word[word.text]
                else:
                    sent_strength[sent]=freq_word[word.text]
    return sent_strength

def summarize(model, file):

    if str(file).endswith('.csv') or str(file).endswith('.xlsx'):
        filename, summary, keywords = ReadFiles2(file)
    else:
        text, filename = ReadFiles(file)

        freq_word, doc, keywords = keyword_extract(model, text)

        sent_strength = sentence_weighting(doc, freq_word)

        summarized_sentences = nlargest(3, sent_strength, key=sent_strength.get)

        final_sentences = [w.text for w in summarized_sentences]
        summary = ' '.join(final_sentences)
        summary = summary.replace("\n", "")

    x = {
        "title": filename,
        "summary": summary,
        "keywords": keywords,
        "subject": "TBD"
    }

    keywords = ", ".join(keywords)
    

    return filename, summary, keywords, json.dumps(x, indent=4)

model = spacy.load('en_core_web_sm')

# filename, summary, keywords, json_result = summarize(model, 'test.pdf')
# print(filename)
# print(summary)
# print(keywords)
# print(json_result['title'])
# print(json_result)
# print(json_result[0]['title'])
# print(json_result[0]['summary'])
# print(json_result[0]['keywords'])


